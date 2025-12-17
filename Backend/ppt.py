import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import colorsys
import re

class ExcelToPPTGenerator:
    def __init__(self, base_path):
        self.base_path = base_path
        self.input_json_path = os.path.join(base_path, "input", "csv", "input.json")
        self.preview_json_path = os.path.join(base_path, "output", "preview.json")
        self.templates_path = os.path.join(base_path, "input", "templates")
        self.cleaned_templates_path = os.path.join(base_path, "input", "cleaned_templates")
        self.default_templates_path = os.path.join(base_path, "default_templates")
        
        # Create necessary directories
        os.makedirs(self.cleaned_templates_path, exist_ok=True)
        os.makedirs(os.path.join(base_path, "output"), exist_ok=True)
    
    def load_json(self, filepath):
        """Load JSON file"""
        with open(filepath, 'r', encoding='utf-8') as f:
            return json.load(f)
    
    def hex_to_rgb(self, hex_color):
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def get_contrasting_color(self, bg_color_rgb):
        """Generate a contrasting color based on background color"""
        r, g, b = [x / 255.0 for x in bg_color_rgb]
        h, s, v = colorsys.rgb_to_hsv(r, g, b)
        
        # If background is dark, return light color; if light, return dark color
        if v < 0.5:
            new_v = 0.9
        else:
            new_v = 0.1
        
        new_r, new_g, new_b = colorsys.hsv_to_rgb(h, s, new_v)
        return (int(new_r * 255), int(new_g * 255), int(new_b * 255))
    
    def extract_slide_background_color(self, slide):
        """Extract background color from slide"""
        try:
            if slide.background.fill.type == 1:  # Solid fill
                color = slide.background.fill.fore_color.rgb
                return (color[0], color[1], color[2])
        except:
            pass
        
        # Default white background if cannot extract
        return (255, 255, 255)
    
    def clean_template(self, template_path, cleaned_path):
        """Remove only text content from text boxes and placeholders, keep all design elements"""
        prs = Presentation(template_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                # Only clear text content, don't remove the shape itself
                if hasattr(shape, "text_frame"):
                    try:
                        # Clear text from the text frame
                        text_frame = shape.text_frame
                        text_frame.clear()
                    except Exception as e:
                        print(f"Could not clear text from shape: {e}")
                
                # If it's a placeholder, clear its text but keep the placeholder
                if shape.is_placeholder:
                    try:
                        if hasattr(shape, "text_frame"):
                            shape.text_frame.clear()
                    except Exception as e:
                        print(f"Could not clear placeholder text: {e}")
        
        prs.save(cleaned_path)
        print(f"Template cleaned - cleared text content, kept all design elements")
        return len(prs.slides)
    
    def calculate_dynamic_font_size(self, text, allocated_width_inches, allocated_height_inches, has_image=False):
        """Calculate optimal font size based on text length and allocated space with better adaptation"""
        
        # Handle different content types
        if isinstance(text, list):
            text_length = sum(len(str(item)) for item in text)
        else:
            text_length = len(str(text))
        
        # Calculate area available in square inches
        area = allocated_width_inches * allocated_height_inches
        
        if area <= 0:
            return 12  # Default fallback
        
        # Calculate text density (characters per square inch)
        chars_per_sq_inch = text_length / area
        
        # More granular font size calculation based on density
        if text_length < 50:
            # Very sparse text - can use large font
            font_size = 34
        elif text_length < 100:
            # Sparse text
            font_size = 32
        elif text_length < 150:
            # Comfortable density
            font_size = 30
        elif text_length < 250:
            # Moderate density
            font_size = 28
        elif text_length < 350:
            # Dense text
            font_size = 26
        elif text_length < 500:
            # Very dense text
            font_size = 24
        elif text_length < 700:
            # Extremely dense text
            font_size = 22
        elif text_length < 900:
            # Extremely dense text
            font_size = 20
        elif text_length < 1200:
            # Extremely dense text
            font_size = 18
        else:
            # Overly dense - minimum readable size
            font_size = 16
        
        # Adjust based on available height
        if allocated_height_inches < 1.0:
            font_size = min(font_size, 12)
        elif allocated_height_inches < 1.5:
            font_size = min(font_size, 14)
        elif allocated_height_inches < 2.0:
            font_size = min(font_size, 16)
        elif allocated_height_inches < 2.5:
            font_size = min(font_size, 18)
        elif allocated_height_inches < 3.0:
            font_size = min(font_size, 20)
        
        # Adjust based on available width (for very narrow columns)
        if allocated_width_inches < 2.0:
            font_size = min(font_size, 14)
        elif allocated_width_inches < 3.0:
            font_size = min(font_size, 16)
        
        # Reduce font size if there's an image sharing the space
        if has_image:
            font_size = max(10, font_size - 2)  # Smaller reduction
        
        # Ensure font size is within reasonable bounds
        font_size = max(10, min(font_size, 36))
        
        print(f"Text adaptation: {text_length} chars, {area:.1f} sq.in, density: {chars_per_sq_inch:.1f} chars/sq.in -> font: {font_size}pt")
        
        return font_size
    
    def add_title_to_slide(self, slide, title, font_color_rgb, slide_width, slide_height, is_first_slide=False, is_last_slide=False):
        """Add title to slide header with improved alignment"""
        if is_first_slide:
            # Center of entire slide for first slide with better vertical centering
            left = Inches(0.5)
            top = slide_height / 2 - Inches(1.2)  # Better centered position
            width = slide_width - Inches(1)
            height = Inches(1.5)
        elif is_last_slide:
            # Center of entire slide for last slide
            left = Inches(0.5)
            top = slide_height / 2 - Inches(1.0)  # Better position for title with subtitle
            width = slide_width - Inches(1)
            height = Inches(1.2)
        else:
            # Header section for other slides with better spacing
            left = Inches(0.5)
            top = Inches(0.4)  # Slightly more top margin
            width = slide_width - Inches(1)
            height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = title
        
        # Improved text formatting
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.line_spacing = 1.2  # Better line spacing
        
        run = paragraph.runs[0]
        if is_first_slide:
            run.font.size = Pt(36)  # Slightly adjusted for better proportion
        elif is_last_slide:
            run.font.size = Pt(36)
        else:
            run.font.size = Pt(28)  # Reduced for better hierarchy
        run.font.bold = True
        run.font.color.rgb = RGBColor(*font_color_rgb)
        
        text_frame.word_wrap = True
        text_frame.auto_size = True  # Disable auto-size for better control
        if is_first_slide or is_last_slide:
            text_frame.vertical_anchor = 1  # Middle alignment
    
    def add_subtitle_to_first_slide(self, slide, subtitle, font_color_rgb, slide_width, slide_height):
        """Add subtitle to first slide at right side from center"""
        # Position subtitle below and to the right of center
        left = slide_width / 2 + Inches(0.5)  # Right side from center
        top = slide_height / 2 + Inches(0.3)   # Below title
        width = slide_width / 2 - Inches(1)    # Half width from center to right edge
        height = Inches(0.6)
        
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = subtitle_box.text_frame
        text_frame.text = subtitle
        
        # Format subtitle
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        
        run = paragraph.runs[0]
        run.font.size = Pt(18)
        run.font.italic = True
        run.font.color.rgb = RGBColor(*font_color_rgb)
        
        text_frame.word_wrap = True
    
    def add_subtitle_to_last_slide(self, slide, subtitle, font_color_rgb, slide_width, slide_height):
        """Add subtitle to last slide centered below title"""
        # Position subtitle centered below title
        left = Inches(0.5)
        top = slide_height / 2 + Inches(0.2)   # Below the centered title
        width = slide_width - Inches(1)
        height = Inches(0.6)
        
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = subtitle_box.text_frame
        text_frame.text = subtitle
        
        # Format subtitle
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        
        run = paragraph.runs[0]
        run.font.size = Pt(20)
        run.font.italic = True
        run.font.color.rgb = RGBColor(*font_color_rgb)
        
        text_frame.word_wrap = True
        text_frame.vertical_anchor = 1  # Middle alignment
    
    def add_content_to_slide(self, slide, content, font_color_rgb, left, top, width, height, has_image=False):
        """Add content to slide body with improved dynamic font sizing"""
        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        
        # Convert content to string if it's a list
        if isinstance(content, list):
            content_text = '\n'.join(str(item) for item in content)
        else:
            content_text = str(content)
        
        text_frame.text = content_text
        
        # Calculate dynamic font size based on actual content and allocated space
        width_inches = width / Inches(1)
        height_inches = height / Inches(1)
        font_size = self.calculate_dynamic_font_size(content_text, width_inches, height_inches, has_image)
        
        # Improved text formatting with adaptive spacing
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.line_spacing = 1.2 + (font_size / 100)  # Adaptive line spacing
            
            # Adaptive paragraph spacing based on font size
            if font_size <= 14:
                paragraph.space_before = Pt(3)
                paragraph.space_after = Pt(1)
            elif font_size <= 18:
                paragraph.space_before = Pt(4)
                paragraph.space_after = Pt(2)
            else:
                paragraph.space_before = Pt(6)
                paragraph.space_after = Pt(3)
            
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.color.rgb = RGBColor(*font_color_rgb)
        
        text_frame.word_wrap = True
        
        # Adaptive margins based on font size
        margin_base = font_size * 0.1
        text_frame.margin_left = Inches(min(0.15, margin_base / 100))
        text_frame.margin_right = Inches(min(0.15, margin_base / 100))
        text_frame.margin_top = Inches(min(0.05, margin_base / 200))
        text_frame.margin_bottom = Inches(min(0.05, margin_base / 200))
        
        return font_size
    
    def add_image_to_slide(self, slide, image_path, left, top, width, height, is_dashboard=False):
        """Add image to slide with dynamic sizing in 2:3 ratio or 3:1 for dashboard"""
        if not os.path.exists(image_path):
            print(f"Warning: Image not found - {image_path}")
            return False
        
        try:
            # Get image dimensions
            img = Image.open(image_path)
            img_width, img_height = img.size
            aspect_ratio = img_width / img_height
            
            # Calculate dimensions to fit within allocated space
            target_width = width
            target_height = target_width / aspect_ratio
            
            # Check if the resulting height exceeds available space
            if target_height > height:
                target_height = height
                target_width = target_height * aspect_ratio
            
            # Apply ratio constraint based on whether it's a dashboard
            if is_dashboard:
                desired_ratio = 4/3  # 3:1 ratio for dashboard
            else:
                desired_ratio = 2/3  # 2:3 ratio for other images
            
            current_ratio = target_width / target_height
            
            if current_ratio > desired_ratio:
                # Too wide - adjust height to maintain desired ratio
                target_height = target_width / desired_ratio
                if target_height > height:
                    target_height = height
                    target_width = target_height * desired_ratio
            else:
                # Too tall - adjust width to maintain desired ratio
                target_width = target_height * desired_ratio
                if target_width > width:
                    target_width = width
                    target_height = target_width / desired_ratio
            
            # Center image in allocated space
            left_offset = left + (width - target_width) / 2
            top_offset = top + (height - target_height) / 2
            
            slide.shapes.add_picture(image_path, left_offset, top_offset, target_width, target_height)
            ratio_type = "4:3" if is_dashboard else "2:3"
            print(f"Added image with size: {target_width / Inches(1):.2f}\" x {target_height / Inches(1):.2f}\" (ratio: {ratio_type})")
            return True
        except Exception as e:
            print(f"Error adding image {image_path}: {e}")
            return False
    
    def add_multiple_images_to_slide(self, slide, image_paths, left, top, width, height, is_dashboard=False):
        """Add multiple images in grid format with improved spacing"""
        valid_images = []
        
        # Check which images exist
        for img_path in image_paths:
            if os.path.exists(img_path):
                valid_images.append(img_path)
            else:
                print(f"Warning: Image not found - {img_path}")
        
        if not valid_images:
            print("No valid images found to display")
            return
        
        num_images = len(valid_images)
        
        # Improved grid layout determination
        if num_images == 1:
            rows, cols = 1, 1
        elif num_images == 2:
            rows, cols = 1, 2
        elif num_images <= 4:
            rows, cols = 2, 2
        elif num_images <= 6:
            rows, cols = 2, 3
        else:
            rows, cols = 3, 3
        
        # Calculate cell dimensions with better spacing
        cell_width = width / cols
        cell_height = height / rows
        
        # Increased padding for better visual separation
        horizontal_padding = Inches(0.15)
        vertical_padding = Inches(0.1)
        
        img_width = cell_width - horizontal_padding * 2
        img_height = cell_height - vertical_padding * 2
        
        for idx, img_path in enumerate(valid_images[:rows * cols]):
            row = idx // cols
            col = idx % cols
            
            # Center images better within their cells
            img_left = left + col * cell_width + horizontal_padding
            img_top = top + row * cell_height + vertical_padding
            
            self.add_image_to_slide(slide, img_path, img_left, img_top, img_width, img_height, is_dashboard)
    
    def extract_image_paths_from_content(self, content_text):
        """Extract image paths from content text"""
        image_paths = []
        
        # Handle both string and list content types
        if isinstance(content_text, list):
            # If content is a list, join it into a string
            content_string = ' '.join(str(item) for item in content_text)
        else:
            # If content is already a string, use it as is
            content_string = str(content_text)
        
        # Pattern to find image paths in the format (Chart: path/to/image.png)
        img_pattern = r'\(Chart:\s*([^)]+)\)'
        matches = re.findall(img_pattern, content_string)
        
        for match in matches:
            # Clean up the path - remove any extra spaces
            img_relative = match.strip()
            
            # Normalize path separators
            img_relative_normalized = img_relative.replace('\\', os.sep).replace('/', os.sep)
            
            # Check if path is relative to base_path
            if not os.path.isabs(img_relative_normalized):
                full_path = os.path.join(self.base_path, img_relative_normalized)
            else:
                full_path = img_relative_normalized
            
            # Check if file exists
            if os.path.exists(full_path):
                if full_path not in image_paths:
                    image_paths.append(full_path)
            else:
                print(f"Warning: Image not found at path: {full_path}")
        
        # Remove image references from content text
        clean_content = re.sub(img_pattern, '', content_string)
        clean_content = clean_content.strip()
        
        return image_paths, clean_content
    
    def get_template_path(self, template_name):
        """Get the correct template path - custom templates have priority, then default templates"""
        # First check if it's a custom template in input/templates folder
        custom_template_path = os.path.join(self.templates_path, template_name)
        if os.path.exists(custom_template_path):
            print(f"Using custom template: {template_name}")
            return custom_template_path
        
        # If not found as custom template, check if it's a default template
        default_templates = ['Creative', 'Professional', 'Minimal', 'Technical']
        
        # Check both exact match and case-insensitive match
        template_lower = template_name.lower()
        for default_template in default_templates:
            if template_lower == default_template.lower():
                default_template_path = os.path.join(self.default_templates_path, f"{default_template}.pptx")
                if os.path.exists(default_template_path):
                    print(f"Using default template: {default_template}")
                    return default_template_path
        
        # If no template found, use Professional as fallback
        fallback_template = "Professional.pptx"
        fallback_path = os.path.join(self.default_templates_path, fallback_template)
        if os.path.exists(fallback_path):
            print(f"Template {template_name} not found, using fallback template: {fallback_template}")
            return fallback_path
        else:
            raise FileNotFoundError(f"Template {template_name} not found and no fallback template available")
    
    def generate_ppt(self):
        """Main function to generate PowerPoint presentation"""
        # Load input files
        input_data = self.load_json(self.input_json_path)
        preview_data = self.load_json(self.preview_json_path)
        
        # Get template info
        presentation_title = input_data['presentation_title']  # Title from input.json for first slide
        template_name = input_data['template_name']
        font_color_hex = input_data['text_color']
        font_color_rgb = self.hex_to_rgb(font_color_hex)
        
        # Get the correct template path
        template_path = self.get_template_path(template_name)
        cleaned_template_path = os.path.join(self.cleaned_templates_path, os.path.basename(template_path))
        
        # Clean template (removes text but keeps design images)
        print(f"Cleaning template: {template_name}")
        num_template_slides = self.clean_template(template_path, cleaned_template_path)
        print(f"Template has {num_template_slides} slides")
        
        # Load cleaned template - DO NOT MODIFY THE ORIGINAL SLIDES
        prs = Presentation(cleaned_template_path)
        
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        print(f"Slide dimensions: {slide_width / Inches(1):.2f}\" x {slide_height / Inches(1):.2f}\"")
        
        # Extract background color from first slide
        if len(prs.slides) > 0:
            bg_color = self.extract_slide_background_color(prs.slides[0])
            print(f"Detected background color: RGB{bg_color}")
            contrasting_color = self.get_contrasting_color(bg_color)
            print(f"Using contrasting color: RGB{contrasting_color}")
        else:
            contrasting_color = font_color_rgb
        
        # Calculate target number of slides
        total_slides_needed = len(preview_data['slides'])
        
        print(f"Total slides needed: {total_slides_needed}")
        print(f"Template slides available: {num_template_slides}")
        
        # Adjust number of slides WITHOUT removing designs
        if num_template_slides >= total_slides_needed:
            # We have enough slides, just use what we have
            print(f"Using existing {num_template_slides} template slides")
        else:
            # Need to add more slides by duplicating with correct mapping
            slides_to_add = total_slides_needed - num_template_slides
            print(f"Adding {slides_to_add} additional slides...")
            
            # Store original slides for reference
            original_slides = list(prs.slides)
            
            # Add slides with correct mapping pattern
            for i in range(slides_to_add):
                ppt_slide_index = num_template_slides + i
                
                # Determine which template slide to use based on mapping pattern
                if ppt_slide_index == 0:
                    template_slide_index = 0  # First slide
                elif ppt_slide_index == total_slides_needed - 1:
                    template_slide_index = num_template_slides - 1  # Last slide
                else:
                    # Middle slides cycle through template slides 1 to (num_template_slides-2)
                    if num_template_slides > 2:
                        template_slide_index = (ppt_slide_index - 1) % (num_template_slides - 2) + 1
                    else:
                        template_slide_index = min(ppt_slide_index, num_template_slides - 1)
                
                # Add the mapped slide
                if template_slide_index < len(original_slides):
                    template_slide = original_slides[template_slide_index]
                    new_slide = prs.slides.add_slide(template_slide.slide_layout)
                    
                    # Copy the design elements from the template slide
                    # This preserves backgrounds, images, and other design elements
                    for shape in template_slide.shapes:
                        if hasattr(shape, 'element'):
                            new_slide.shapes._spTree.append(shape.element)
                    
                    print(f"PPT Slide {ppt_slide_index + 1} -> Template Slide {template_slide_index + 1}")
                else:
                    # Fallback
                    new_slide = prs.slides.add_slide(prs.slide_layouts[0])
                    print(f"PPT Slide {ppt_slide_index + 1} -> Default Layout")
        
        print(f"\nGenerating {total_slides_needed} slides with preserved designs...")
        
        # Process each slide from preview.json
        for slide_data in preview_data['slides']:
            slide_index = slide_data['slide_index'] - 1
            if slide_index >= len(prs.slides):
                print(f"Warning: Slide index {slide_index + 1} exceeds available slides")
                continue
                
            slide = prs.slides[slide_index]
            placeholders = slide_data['placeholders']
            
            # Determine slide type
            is_first_slide = slide_index == 0
            is_last_slide = slide_index == total_slides_needed - 1
            
            # Check if this is a dashboard slide
            title_text = placeholders.get('title', '').lower() if not is_first_slide else ''
            is_dashboard_slide = 'dashboard' in title_text
            
            # Calculate body area with better margins
            body_top = Inches(1.5)  # Increased from 1.2 for better title-body separation
            body_height = slide_height - Inches(2.0)  # Increased bottom margin
            body_width = slide_width - Inches(1.2)    # Increased side margins
            
            # For first and last slides, adjust body area
            if is_first_slide or is_last_slide:
                body_top = slide_height / 2 + Inches(0.5)  # Better positioning for centered slides
                body_height = slide_height / 2 - Inches(1.0)
            
            # Add title
            if is_first_slide:
                # Use title from input.json for first slide
                title_text = presentation_title
            else:
                # Use title from preview.json for other slides
                title_text = placeholders.get('title', '')
            
            if title_text:
                self.add_title_to_slide(slide, title_text, font_color_rgb, slide_width, slide_height, is_first_slide, is_last_slide)
            
            # Handle first slide subtitle
            if is_first_slide and 'subtitle' in placeholders:
                subtitle = placeholders['subtitle']
                self.add_subtitle_to_first_slide(slide, subtitle, font_color_rgb, slide_width, slide_height)
                continue
            
            # Handle last slide subtitle
            if is_last_slide and 'subtitle' in placeholders:
                subtitle = placeholders['subtitle']
                self.add_subtitle_to_last_slide(slide, subtitle, font_color_rgb, slide_width, slide_height)
                continue
            
            # Skip body for last slide if only title
            if is_last_slide and 'content' not in placeholders:
                continue
            
            # Get content
            content_text = placeholders.get('content', '')
            
            # Check for image_path in placeholders (from preview.json)
            image_path = placeholders.get('image_path', '')
            
            # For dashboard slides, prioritize the image_path and use full slide for image
            if is_dashboard_slide and image_path and os.path.exists(image_path):
                # Dashboard slide - use full slide for image with 3:1 ratio
                dashboard_left = Inches(0.5)
                dashboard_top = Inches(1.2)
                dashboard_width = slide_width - Inches(1)
                dashboard_height = slide_height - Inches(2.0)
                
                image_added = self.add_image_to_slide(slide, image_path, 
                                      dashboard_left, dashboard_top, dashboard_width, dashboard_height, True)
                if not image_added:
                    print(f"Slide {slide_index + 1}: Dashboard image not found at {image_path}")
                continue
            
            if not content_text:
                continue
            
            if image_path and os.path.exists(image_path):
                # Single image from image_path
                image_paths = [image_path]
                clean_content = content_text
            else:
                # Extract image paths from content
                image_paths, clean_content = self.extract_image_paths_from_content(content_text)
            
            # Get content length for dynamic layout
            if isinstance(clean_content, list):
                content_length = sum(len(str(item)) for item in clean_content)
            else:
                content_length = len(str(clean_content))
            
            # Layout based on content and images with dynamic ratios
            if image_paths and clean_content:
                # Determine optimal split ratio based on content length
                if content_length < 100:
                    # Short content - give more space to image
                    content_width = body_width * 0.4
                    image_width = body_width * 0.55
                elif content_length < 300:
                    # Medium content - balanced split
                    content_width = body_width * 0.5
                    image_width = body_width * 0.45
                else:
                    # Long content - give more space to text
                    content_width = body_width * 0.6
                    image_width = body_width * 0.35
                
                gap = Inches(0.3)
                
                # Add content on left with better spacing
                content_left = Inches(0.5)
                content_top_adjusted = body_top + Inches(0.1)  # Slight adjustment
                font_size = self.add_content_to_slide(slide, clean_content, font_color_rgb, 
                                    content_left, content_top_adjusted, content_width, body_height - Inches(0.1), has_image=True)
                
                # Add images on right with better positioning
                image_left = Inches(0.5) + content_width + gap
                image_top_adjusted = body_top + Inches(0.1)
                if len(image_paths) == 1:
                    image_added = self.add_image_to_slide(slide, image_paths[0], 
                                          image_left, image_top_adjusted, image_width, body_height - Inches(0.1), is_dashboard_slide)
                    if not image_added:
                        print(f"Slide {slide_index + 1}: Image space left empty (image not found)")
                else:
                    self.add_multiple_images_to_slide(slide, image_paths, 
                                                 image_left, image_top_adjusted, image_width, body_height - Inches(0.1), is_dashboard_slide)
                
                print(f"Slide {slide_index + 1}: Content length {content_length} chars -> font: {font_size}pt, layout: {content_width/body_width*100:.0f}% text, {image_width/body_width*100:.0f}% image")
            
            elif image_paths and not clean_content:
                # Only images, no content
                image_left = Inches(0.5)
                image_top_adjusted = body_top + Inches(0.1)
                if len(image_paths) == 1:
                    self.add_image_to_slide(slide, image_paths[0], 
                                          image_left, image_top_adjusted, body_width, body_height - Inches(0.1), is_dashboard_slide)
                else:
                    self.add_multiple_images_to_slide(slide, image_paths, 
                                                 image_left, image_top_adjusted, body_width, body_height - Inches(0.1), is_dashboard_slide)
                print(f"Slide {slide_index + 1}: Images only layout")
            
            elif clean_content and not image_paths:
                # Only content, no images
                content_left = Inches(0.5)
                content_top_adjusted = body_top + Inches(0.1)
                font_size = self.add_content_to_slide(slide, clean_content, font_color_rgb, 
                                    content_left, content_top_adjusted, body_width, body_height - Inches(0.1), has_image=False)
                print(f"Slide {slide_index + 1}: Content length {content_length} chars -> font: {font_size}pt (full width)")
        
        # Save presentation
        output_path = os.path.join(self.base_path, "output", "generated_presentation.pptx")
        prs.save(output_path)
        print(f"\nPresentation saved to: {output_path}")
        return output_path


# Main execution
if __name__ == "__main__":
    # Set your base path
    base_path = r"C:\Users\raj.kumar\Documents\excel-ppt project\backend"  # Update this path    
    generator = ExcelToPPTGenerator(base_path)
    output_file = generator.generate_ppt()
    print(f"Successfully generated: {output_file}")  