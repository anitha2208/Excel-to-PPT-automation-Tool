from fastapi import FastAPI, File, UploadFile, HTTPException, BackgroundTasks, Form
from fastapi.responses import FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
import pandas as pd
import os
import uuid
import shutil
from pathlib import Path
import json
from datetime import datetime
from typing import Dict, List, Optional, Any
import asyncio
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import logging
import matplotlib.pyplot as plt
import seaborn as sns
from io import BytesIO
import base64
import subprocess
import sys




# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="Excel to PPT Generator API", version="1.0.0")

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configuration
BASE_DIR = Path(__file__).parent
INPUT_DIR = BASE_DIR / "input"
OUTPUT_DIR = BASE_DIR / "output"

# Create input subdirectories
CSV_DIR = INPUT_DIR / "csv"
TEMPLATE_DIR = INPUT_DIR / "templates"
LOGO_DIR = INPUT_DIR / "logos"

# Create directories
for directory in [INPUT_DIR, OUTPUT_DIR, CSV_DIR, TEMPLATE_DIR, LOGO_DIR]:
    directory.mkdir(parents=True, exist_ok=True)
    logger.info(f"Created directory: {directory}")

# Mount static directories for file serving
app.mount("/input", StaticFiles(directory=INPUT_DIR), name="input")
app.mount("/output", StaticFiles(directory=OUTPUT_DIR), name="output")

# Pydantic models
class AnalysisRequest(BaseModel):
    file_id: str

class ChartGenerationRequest(BaseModel):
    column: str
    chart_type: str

class PPTGenerationRequest(BaseModel):
    csv_file_id: str
    template_file_id: Optional[str] = None
    logo_file_id: Optional[str] = None
    presentation_config: Dict[str, Any] = {}
    options: Dict[str, Any] = {}

class FileUploadResponse(BaseModel):
    file_id: str
    filename: str
    file_type: str
    file_path: str
    message: str

class AnalysisResponse(BaseModel):
    file_id: str
    columns: List[str]
    summary: Dict[str, Any]
    sample_data: List[Dict]

class SlidePreviewResponse(BaseModel):
    success: bool
    slide_preview_data: List[Dict]
    total_slides: int
    timestamp: str

class TemplateResponse(BaseModel):
    success: bool
    templates: List[Dict]
    current_template: str

# NEW: Slide Images Response Model
class SlideImagesResponse(BaseModel):
    success: bool
    ppt_id: str
    slide_images: List[Dict]
    total_slides: int
    timestamp: str

# Utility functions
def cleanup_input_output_folders():
    """Completely clean input and output folders when new dataset is loaded"""
    try:
        logger.info("Cleaning up input and output folders...")
        
        # Clean OUTPUT_DIR completely
        if OUTPUT_DIR.exists():
            shutil.rmtree(OUTPUT_DIR)
            OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
            logger.info(f"Cleaned output directory: {OUTPUT_DIR}")
        
        # Clean INPUT_DIR completely but recreate necessary subdirectories
        if INPUT_DIR.exists():
            shutil.rmtree(INPUT_DIR)
            INPUT_DIR.mkdir(parents=True, exist_ok=True)
            
            # Recreate subdirectories
            for directory in [CSV_DIR, TEMPLATE_DIR, LOGO_DIR]:
                directory.mkdir(parents=True, exist_ok=True)
                logger.info(f"Recreated directory: {directory}")
        
        logger.info("Input and output folders cleaned successfully")
        return True
        
    except Exception as e:
        logger.error(f"Error cleaning input/output folders: {e}")
        return False

def get_file_path(file_id: str, file_type: str = "csv") -> Path:
    """Get file path from file_id and type"""
    if file_type == "template":
        return TEMPLATE_DIR / f"{file_id}.pptx"
    elif file_type == "logo":
        # Logo can have various extensions, so we need to find the actual file
        logo_dir = LOGO_DIR
        for file in logo_dir.glob(f"{file_id}.*"):
            return file
        return LOGO_DIR / f"{file_id}"  # Fallback
    else:
        # For CSV/Excel files, check for both extensions
        csv_path = CSV_DIR / "input.csv"
        excel_path = CSV_DIR / "input.xlsx"
        
        if excel_path.exists():
            return excel_path
        elif csv_path.exists():
            return csv_path
        else:
            return csv_path  # fallback

def save_uploaded_file(file: UploadFile, file_type: str = "csv") -> tuple[str, str]:
    """Save uploaded file and return (file_id, filename) - NO SIZE LIMIT"""
    try:
        original_filename = file.filename or "uploaded_file"
        file_extension = Path(original_filename).suffix.lower()
        
        # Clean up folders when new CSV dataset is loaded
        if file_type == "csv":
            cleanup_input_output_folders()
            
        # Generate fixed file_id for CSV files
        if file_type == "csv":
            file_id = "input"  # Fixed file_id for CSV files
            save_dir = CSV_DIR
            
            # Determine filename based on file type
            if file_extension in ['.xlsx', '.xls', '.xlsv']:
                filename = "input.xlsx"
            else:
                filename = "input.csv"

        elif file_type == "template":
            original_filename = file.filename or "template.pptx"
            save_dir = TEMPLATE_DIR
            filename = original_filename
            file_id = str(uuid.uuid4())
        
        file_path = save_dir / filename
        
        logger.info(f"Saving {file_type} file: {file_path}")
        
        # Remove existing file if it exists
        if file_path.exists():
            file_path.unlink()
            logger.info(f"Removed existing file: {file_path}")
        
        # Save the file - NO SIZE LIMIT
        with open(file_path, "wb") as buffer:
            # Read file in chunks to handle very large files
            file_size = 0
            chunk_size = 1024 * 1024  # 1MB chunks
            
            while True:
                chunk = file.file.read(chunk_size)
                if not chunk:
                    break
                buffer.write(chunk)
                file_size += len(chunk)
            
            file_size_mb = file_size / (1024 * 1024)
            logger.info(f"File saved successfully: {file_path} (Size: {file_size_mb:.2f} MB)")
        
        return file_id, filename
        
    except Exception as e:
        logger.error(f"Error saving file: {e}")
        raise HTTPException(status_code=500, detail=f"File save failed: {str(e)}")

def analyze_csv_data(file_path: Path) -> Dict[str, Any]:
    """Analyze CSV data and return insights with large file support - NO SIZE LIMIT"""
    try:
        logger.info(f"Analyzing file: {file_path}")
        
        # Get file size for optimization
        file_size = file_path.stat().st_size / (1024 * 1024)  # Size in MB
        logger.info(f"File size: {file_size:.2f} MB")
        
        # Memory optimization for large files - NO SIZE LIMIT, but use sampling for very large files
        sample_size = 100000  # Increased sample size for better analysis
        
        # Read the file based on extension with optimization for large files
        if file_path.suffix.lower() in ['.xlsx', '.xls', '.xlsv']:
            # For large Excel files, use sampling if very large
            if file_size > 100:  # > 100MB use sampling
                logger.info(f"Large Excel file detected, using sample of {sample_size} rows")
                df = pd.read_excel(file_path, nrows=sample_size)
            else:
                df = pd.read_excel(file_path)
        else:
            # For large CSV files, use sampling if very large
            if file_size > 100:  # > 100MB use sampling
                logger.info(f"Large CSV file detected, using sample of {sample_size} rows")
                try:
                    df = pd.read_csv(file_path, nrows=sample_size)
                except UnicodeDecodeError:
                    df = pd.read_csv(file_path, encoding='latin-1', nrows=sample_size)
            else:
                try:
                    df = pd.read_csv(file_path)
                except UnicodeDecodeError:
                    df = pd.read_csv(file_path, encoding='latin-1')
        
        # Basic analysis
        analysis = {
            "columns": df.columns.tolist(),
            "shape": df.shape,
            "data_types": df.dtypes.astype(str).to_dict(),
            "file_size_mb": round(file_size, 2),
            "is_sampled": file_size > 100,  # Flag if sampling was used
            "sample_size": sample_size if file_size > 100 else df.shape[0],
            "summary": {
                "numeric_columns": df.select_dtypes(include=['number']).columns.tolist(),
                "categorical_columns": df.select_dtypes(include=['object']).columns.tolist(),
                "missing_values": df.isnull().sum().to_dict(),
                "basic_stats": df.describe().to_dict() if not df.select_dtypes(include=['number']).empty else {}
            },
            "sample_data": df.replace({pd.NaT: None, pd.NA: None}).fillna('').to_dict('records')  # Limit sample data for frontend
        }
        
        logger.info(f"Analysis completed. Columns: {len(analysis['columns'])}, Rows: {analysis['shape'][0]}")
        return analysis
        
    except Exception as e:
        logger.error(f"Error analyzing file: {e}")
        raise HTTPException(status_code=500, detail=f"Error analyzing file: {str(e)}")

def generate_summary_statistics(df: pd.DataFrame) -> dict:
    """Generate summary statistics for numeric columns"""
    try:
        numeric_df = df.select_dtypes(include=['number'])
        if numeric_df.empty:
            return {}
        
        stats = numeric_df.describe().to_dict()
        
        # Add additional statistics
        for col in numeric_df.columns:
            stats[col].update({
                'variance': float(numeric_df[col].var()) if not pd.isna(numeric_df[col].var()) else 0.0,
                'skewness': float(numeric_df[col].skew()) if not pd.isna(numeric_df[col].skew()) else 0.0,
                'kurtosis': float(numeric_df[col].kurtosis()) if not pd.isna(numeric_df[col].kurtosis()) else 0.0,
                'median': float(numeric_df[col].median()) if not pd.isna(numeric_df[col].median()) else 0.0,
            })
            # Handle mode carefully
            mode_values = numeric_df[col].mode()
            if not mode_values.empty:
                stats[col]['mode'] = float(mode_values.iloc[0])
            else:
                stats[col]['mode'] = None
        
        return stats
    except Exception as e:
        logger.error(f"Error generating summary statistics: {e}")
        return {}

def generate_slide_preview_data(config: dict, df: pd.DataFrame = None) -> List[Dict]:
    """Generate slide preview data for frontend display"""
    try:
        slide_data = []
        features = config.get("features", {})
        
        # Title Slide
        slide_data.append({
            "title": "📋 Title Slide",
            "content": [
                f"**Title:** {config.get('presentation_title', 'Data Analysis Presentation')}",
                f"**Template:** {config.get('template_name', 'Professional').title()}",
                f"**Text Color:** {config.get('text_color', '#ffffff')}",
                f"**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            ],
            "slide_number": 1,
            "type": "title"
        })
        
        # Data Overview Slide
        if df is not None:
            slide_data.append({
                "title": "📊 Data Overview",
                "content": [
                    f"**Total Rows:** {len(df):,}",
                    f"**Total Columns:** {len(df.columns)}",
                    f"**Numeric Columns:** {len(df.select_dtypes(include=['number']).columns)}",
                    f"**Categorical Columns:** {len(df.select_dtypes(include=['object']).columns)}"
                ],
                "slide_number": 2,
                "type": "overview"
            })
        
        slide_counter = 3
        
        # Dashboard Slides
        if features.get("dashboard"):
            dashboard_content = ["**Dashboard Charts:**"]
            for col, chart_config in features["dashboard"].items():
                if isinstance(chart_config, dict):
                    chart_type = chart_config.get("type", "chart")
                    y_column = chart_config.get("y_column", "None")
                    if y_column != "None":
                        dashboard_content.append(f"• {chart_config['x_column']} vs {y_column}: {chart_type} chart")
                    else:
                        dashboard_content.append(f"• {chart_config['x_column']}: {chart_type} chart")
                else:
                    dashboard_content.append(f"• {col}: {chart_config} chart")
            
            slide_data.append({
                "title": "📈 Dashboard Overview",
                "content": dashboard_content,
                "slide_number": slide_counter,
                "type": "dashboard"
            })
            slide_counter += 1
        
        # Insights Slides
        if features.get("insights"):
            insights_content = ["**AI Insights for:**"]
            for col in features["insights"]:
                insights_content.append(f"• {col}")
            
            slide_data.append({
                "title": "💡 AI Insights",
                "content": insights_content,
                "slide_number": slide_counter,
                "type": "insights"
            })
            slide_counter += 1
        
        # Charts with Insights Slides
        if features.get("insights_charts"):
            charts_content = ["**Charts with Insights:**"]
            for col, chart_config in features["insights_charts"].items():
                if isinstance(chart_config, dict):
                    chart_type = chart_config.get("type", "chart")
                    y_column = chart_config.get("y_column", "None")
                    if y_column != "None":
                        charts_content.append(f"• {chart_config['x_column']} vs {y_column}: {chart_type} chart")
                    else:
                        charts_content.append(f"• {chart_config['x_column']}: {chart_type} chart")
                else:
                    charts_content.append(f"• {col}: {chart_config} chart")
            
            slide_data.append({
                "title": "📊 Charts with Insights",
                "content": charts_content,
                "slide_number": slide_counter,
                "type": "insights_charts"
            })
            slide_counter += 1
        
        # Comparison Slides
        if features.get("comparison"):
            comparison_content = ["**Data Comparisons:**"]
            for comp, chart_type in features["comparison"].items():
                comparison_content.append(f"• {comp}: {chart_type} chart")
            
            slide_data.append({
                "title": "⚖️ Data Comparisons",
                "content": comparison_content,
                "slide_number": slide_counter,
                "type": "comparison"
            })
            slide_counter += 1
        
        # Query Slides
        if features.get("queries"):
            query_content = ["**Custom Queries:**"]
            display_as_table = features.get("display_queries_as_table", False)
            
            for i, query in enumerate(features["queries"]):
                query_data = query if isinstance(query, dict) else {"text": str(query)}
                query_text = f"• {query_data.get('text', 'Unknown query')}"
                
                if display_as_table or query_data.get('display_as_table', False):
                    query_text += " (Table format)"
                else:
                    query_text += " (Chart format)"
                
                if not query_data.get('include_in_slides', True):
                    query_text += " [EXCLUDED]"
                
                query_content.append(query_text)
            
            slide_data.append({
                "title": "🔍 Custom Queries",
                "content": query_content,
                "slide_number": slide_counter,
                "type": "queries"
            })
            slide_counter += 1
        
        # Summary Slide
        summary_content = [
            "**Presentation Summary:**",
            f"• Total Slides: {slide_counter}",
            f"• Data Source: {config.get('data_source', 'Uploaded Dataset')}",
            f"• Template: {config.get('template_name', 'Professional').title()}",
            f"• Text Color: {config.get('text_color', '#ffffff')}",
            "• Analysis Complete: Yes",
            f"• Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            "",
            "**Features Included:**"
        ]
        
        if features.get("dashboard"):
            summary_content.append("• Dashboard Charts")
        if features.get("insights"):
            summary_content.append("• AI Insights")
        if features.get("insights_charts"):
            summary_content.append("• Charts with Insights")
        if features.get("comparison"):
            summary_content.append("• Data Comparisons")
        if features.get("queries"):
            summary_content.append("• Custom Queries")
        
        slide_data.append({
            "title": "📝 Summary",
            "content": summary_content,
            "slide_number": slide_counter,
            "type": "summary"
        })
        
        return slide_data
        
    except Exception as e:
        logger.error(f"Error generating slide preview data: {e}")
        return []

def hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def apply_text_color(shape, text_color: str):
    """Apply text color to shape"""
    try:
        rgb_color = hex_to_rgb(text_color)
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(*rgb_color)
    except Exception as e:
        logger.warning(f"Could not apply text color: {e}")

def create_chart_slide(prs, title: str, chart_config: dict, df: pd.DataFrame, text_color: str = "#ffffff"):
    """Create a slide with chart configuration"""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide_title = slide.shapes.title
        slide_title.text = title
        apply_text_color(slide_title, text_color)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = f"Chart Configuration:\n"
        
        if isinstance(chart_config, dict):
            chart_type = chart_config.get("type", "chart")
            x_column = chart_config.get("x_column", "")
            y_column = chart_config.get("y_column", "None")
            
            if y_column != "None":
                tf.text += f"• Chart Type: {chart_type}\n• X-Axis: {x_column}\n• Y-Axis: {y_column}"
            else:
                tf.text += f"• Chart Type: {chart_type}\n• Column: {x_column}"
        
        apply_text_color(content, text_color)
        return slide
    except Exception as e:
        logger.error(f"Error creating chart slide: {e}")
        return None

# ENHANCED FUNCTION: Convert PowerPoint slides to high-quality images with ACTUAL CONTENT
def convert_pptx_to_slide_images(pptx_path: Path) -> List[Dict]:
    """Convert PowerPoint slides to base64 encoded images with ACTUAL SLIDE CONTENT"""
    try:
        logger.info(f"Converting PPTX to actual slide images: {pptx_path}")
        
        if not pptx_path.exists():
            logger.error(f"PPTX file not found: {pptx_path}")
            return []
        
        slide_images = []
        prs = Presentation(pptx_path)
        
        for i, slide in enumerate(prs.slides):
            try:
                # Extract ACTUAL slide content
                slide_title = "Untitled Slide"
                slide_content = []
                has_charts = False
                has_images = False
                
                # Enhanced content extraction from actual slide
                for shape in slide.shapes:
                    # Check for charts
                    if hasattr(shape, 'chart'):
                        has_charts = True
                        chart_type = type(shape.chart).__name__
                        slide_content.append(f"📊 Chart: {chart_type}")
                    
                    # Check for images
                    if hasattr(shape, 'image'):
                        has_images = True
                        slide_content.append(f"🖼️ Image: Embedded visual")
                    
                    # Extract text content
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        
                        # Identify title (usually first text shape or centered/large text)
                        if shape == slide.shapes.title or (not slide_content and len(text) < 100):
                            slide_title = text
                        else:
                            # Add meaningful text content
                            if text and text != slide_title:
                                # Clean up text and split long content
                                lines = text.split('\n')
                                for line in lines:
                                    if line.strip() and len(line.strip()) > 2:
                                        slide_content.append(line.strip())
                
                # Limit content for display but preserve important information
                display_content = slide_content[:15]  # Show more content lines
                
                # Create enhanced visualization of ACTUAL slide
                fig, ax = plt.subplots(figsize=(16, 9), dpi=150)
                ax.axis('off')
                ax.set_facecolor('#2d3748')  # Dark background matching frontend
                
                # Title with enhanced styling
                ax.text(0.5, 0.85, slide_title, 
                       ha='center', va='center', fontsize=22, color='white', weight='bold',
                       transform=ax.transAxes, fontfamily='Arial',
                       bbox=dict(boxstyle="round,pad=0.3", facecolor='#f26f21', alpha=0.9))
                
                # Add visual indicators for slide type
                indicators = []
                if has_charts:
                    indicators.append("📊 Charts")
                if has_images:
                    indicators.append("🖼️ Images")
                if indicators:
                    indicator_text = " | ".join(indicators)
                    ax.text(0.5, 0.78, indicator_text, 
                           ha='center', va='center', fontsize=14, color='#ffa800', weight='bold',
                           transform=ax.transAxes, fontfamily='Arial')
                
                # Add actual content lines with better formatting
                y_start = 0.70
                line_height = 0.05
                
                for j, line in enumerate(display_content):
                    y_pos = y_start - (j * line_height)
                    if y_pos > 0.1:  # Ensure we don't go below bottom
                        ax.text(0.08, y_pos, f"• {line}", 
                               ha='left', va='center', fontsize=12, color='#e2e8f0',
                               transform=ax.transAxes, wrap=True, fontfamily='Arial')
                
                # Add slide number with better styling
                ax.text(0.95, 0.05, f"Slide {i+1}", 
                       ha='right', va='bottom', fontsize=12, color='#ffa800', weight='bold',
                       transform=ax.transAxes, fontfamily='Arial',
                       bbox=dict(boxstyle="circle,pad=0.2", facecolor='#4a5568', alpha=0.8))
                
                # Add a subtle border
                rect = plt.Rectangle((0.02, 0.02), 0.96, 0.96, linewidth=2, 
                                   edgecolor='#f26f21', facecolor='none', transform=ax.transAxes)
                ax.add_patch(rect)
                
                # Convert to base64 with higher quality
                buf = BytesIO()
                plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', 
                           facecolor='#2d3748', edgecolor='none', transparent=False,
                           pad_inches=0.1)
                plt.close(fig)
                
                buf.seek(0)
                img_base64 = base64.b64encode(buf.getvalue()).decode('utf-8')
                
                slide_images.append({
                    "slide_number": i + 1,
                    "image_data": f"data:image/png;base64,{img_base64}",
                    "title": slide_title,
                    "content": display_content,
                    "has_charts": has_charts,
                    "has_images": has_images,
                    "total_shapes": len(slide.shapes),
                    "type": "actual_slide"
                })
                
            except Exception as slide_error:
                logger.warning(f"Error processing slide {i+1}: {slide_error}")
                # Create a fallback representation for this slide
                fig, ax = plt.subplots(figsize=(16, 9), dpi=150)
                ax.axis('off')
                ax.set_facecolor('#2d3748')
                
                ax.text(0.5, 0.6, f"Slide {i+1}", 
                       ha='center', va='center', fontsize=28, color='white', weight='bold',
                       transform=ax.transAxes, fontfamily='Arial')
                ax.text(0.5, 0.45, "Content Loading...", 
                       ha='center', va='center', fontsize=18, color='#a0aec0',
                       transform=ax.transAxes, fontfamily='Arial')
                
                # Convert to base64
                buf = BytesIO()
                plt.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                           facecolor='#2d3748', edgecolor='none')
                plt.close(fig)
                
                buf.seek(0)
                img_base64 = base64.b64encode(buf.getvalue()).decode('utf-8')
                
                slide_images.append({
                    "slide_number": i + 1,
                    "image_data": f"data:image/png;base64,{img_base64}",
                    "title": f"Slide {i + 1}",
                    "content": ["Content processing in progress..."],
                    "has_charts": False,
                    "has_images": False,
                    "total_shapes": 0,
                    "type": "fallback_slide"
                })
        
        logger.info(f"Successfully converted {len(slide_images)} ACTUAL slides to high-quality images")
        return slide_images
        
    except Exception as e:
        logger.error(f"Error converting PPTX to actual slide images: {e}")
        return []

def run_generation_pipeline():
    """Run the generation.py pipeline to process data and create presentation"""
    try:
        logger.info("Starting generation pipeline...")
        
        # Get the path to generation.py
        generation_script = BASE_DIR / "generation.py"
        
        if not generation_script.exists():
            logger.error(f"Generation script not found: {generation_script}")
            return False, "Generation script not found"
        
        # Create environment with forced UTF-8 encoding
        env = os.environ.copy()
        env['PYTHONUTF8'] = '1'
        env['PYTHONIOENCODING'] = 'utf-8'
        
        # Run generation.py as a subprocess with proper encoding and extended timeout
        process = subprocess.Popen(
            [sys.executable, str(generation_script)],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            bufsize=1,
            universal_newlines=False,  # Don't use text mode to avoid encoding issues
            cwd=BASE_DIR,
            env=env
        )
        
        # Read output with explicit UTF-8 decoding and error handling
        try:
            stdout, stderr = process.communicate(timeout=1800)  # 30 minute timeout for large files
            
            # Decode with UTF-8 and replace errors
            stdout_decoded = stdout.decode('utf-8', errors='replace') if stdout else ""
            stderr_decoded = stderr.decode('utf-8', errors='replace') if stderr else ""
            
        except subprocess.TimeoutExpired:
            process.kill()
            stdout_decoded, stderr_decoded = process.communicate()
            stdout_decoded = stdout_decoded.decode('utf-8', errors='replace') if stdout_decoded else ""
            stderr_decoded = stderr_decoded.decode('utf-8', errors='replace') if stderr_decoded else ""
            logger.error("Generation pipeline timed out after 30 minutes")
            return False, "Generation pipeline timed out"
        
        # Log the output for debugging
        if stdout_decoded:
            # Log first 500 characters to avoid huge logs
            logger.info(f"Generation pipeline stdout: {stdout_decoded[:500]}...")
        if stderr_decoded:
            logger.error(f"Generation pipeline stderr: {stderr_decoded}")
        
        if process.returncode == 0:
            logger.info("Generation pipeline completed successfully")
            return True, stdout_decoded
        else:
            logger.error(f"Generation pipeline failed with return code {process.returncode}")
            return False, stderr_decoded
            
    except Exception as e:
        logger.error(f"Error running generation pipeline: {e}")
        return False, str(e)

def debug_file_locations():
    """Debug function to find all PPTX files in the project"""
    try:
        pptx_files = []
        for root, dirs, files in os.walk(BASE_DIR):
            for file in files:
                if file.endswith('.pptx'):
                    full_path = Path(root) / file
                    pptx_files.append({
                        'path': str(full_path),
                        'size': full_path.stat().st_size,
                        'modified': datetime.fromtimestamp(full_path.stat().st_mtime)
                    })
        
        # Sort by modification time (newest first)
        pptx_files.sort(key=lambda x: x['modified'], reverse=True)
        return pptx_files
    except Exception as e:
        logger.error(f"Error debugging file locations: {e}")
        return []

def find_generated_pptx():
    """Find the generated PowerPoint file after generation pipeline runs"""
    try:
        # Look for the latest PPTX file in the output directory
        pptx_files = list(OUTPUT_DIR.glob("*.pptx"))
        
        if not pptx_files:
            # Also check for the specific filename that generation.py creates
            specific_pptx = OUTPUT_DIR / "generated_presentation.pptx"
            if specific_pptx.exists():
                logger.info(f"Found generated PPTX: {specific_pptx}")
                return specific_pptx
            
            # Debug: list all PPTX files in project
            all_pptx = debug_file_locations()
            if all_pptx:
                logger.info(f"Found {len(all_pptx)} PPTX files in project:")
                for pptx in all_pptx[:5]:  # Show first 5
                    logger.info(f"  - {pptx['path']} (Size: {pptx['size']} bytes)")
            
            # Check if generation created it with a different name
            all_files = list(OUTPUT_DIR.glob("*"))
            logger.info(f"All files in output directory: {[f.name for f in all_files]}")
        
        if pptx_files:
            # Get the most recently modified file
            latest_pptx = max(pptx_files, key=lambda x: x.stat().st_mtime)
            logger.info(f"Found generated PPTX: {latest_pptx}")
            return latest_pptx
        else:
            logger.warning("No PPTX files found after generation")
            return None
            
    except Exception as e:
        logger.error(f"Error finding generated PPTX: {e}")
        return None

def generate_powerpoint(request: PPTGenerationRequest) -> dict:
    """Generate PowerPoint presentation based on configuration - SUPPORTS LARGE FILES"""
    try:
        logger.info("Starting PowerPoint generation...")
        
        # Get CSV file path
        csv_path = get_file_path(request.csv_file_id, "csv")
        if not csv_path.exists():
            logger.error(f"CSV file not found: {csv_path}")
            raise HTTPException(status_code=404, detail="CSV file not found")
        
        # First run the generation pipeline
        logger.info("Running generation pipeline...")
        success, message = run_generation_pipeline()
        
        if not success:
            logger.error(f"Generation pipeline failed: {message}")
            # Add more specific error message
            if "timed out" in message.lower():
                raise HTTPException(status_code=500, detail="Generation pipeline timed out. The process took too long to complete.")
            else:
                raise HTTPException(status_code=500, detail=f"Generation pipeline failed: {message}")
        
        # Find the generated PPTX file
        generated_pptx = find_generated_pptx()
        
        if not generated_pptx:
            # More detailed error information
            error_msg = "No PowerPoint file generated by the pipeline. "
            error_msg += "This could be due to: 1) Chart generation failures, 2) Memory issues, 3) PPT template problems"
            logger.error(error_msg)
            raise HTTPException(status_code=500, detail=error_msg)
        
        # Generate a unique ID for the presentation
        ppt_id = str(uuid.uuid4())
        
        # Copy the generated file to our output directory with the new ID
        final_output_path = OUTPUT_DIR / f"{ppt_id}.pptx"
        shutil.copy2(str(generated_pptx), str(final_output_path))
        
        logger.info(f"Presentation copied to: {final_output_path}")
        
        # ENHANCED: Generate high-quality ACTUAL slide images for preview
        slide_images = convert_pptx_to_slide_images(final_output_path)
        
        # Read data for slide preview
        if csv_path.suffix.lower() in ['.xlsx', '.xls', '.xlsv']:
            df = pd.read_excel(csv_path)
        else:
            df = pd.read_csv(csv_path)
        
        # Get presentation configuration
        config = request.presentation_config
        
        # Generate slide preview data
        slide_preview_data = generate_slide_preview_data(config, df)
        
        return {
            "ppt_id": ppt_id,
            "total_slides": len(slide_preview_data),
            "file_path": str(final_output_path),
            "slide_preview_data": slide_preview_data,
            "slide_images": slide_images,  # ENHANCED: Include high-quality ACTUAL slide images
            "generated_at": datetime.now().isoformat(),
            "template_used": config.get("template_name", "generated"),
            "text_color_applied": config.get("text_color", "#ffffff"),
            "generation_pipeline": True
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error generating PowerPoint: {e}")
        import traceback
        logger.error(f"Traceback: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Error generating presentation: {str(e)}")

# ENHANCED FUNCTION: Get actual slide preview from generated PPTX with REAL CONTENT
def get_actual_slide_preview(ppt_id: str) -> List[Dict]:
    """Get actual slide preview from generated PowerPoint file with enhanced content extraction"""
    try:
        pptx_path = OUTPUT_DIR / f"{ppt_id}.pptx"
        if not pptx_path.exists():
            logger.error(f"PPTX file not found for preview: {pptx_path}")
            return []
        
        slide_preview_data = []
        prs = Presentation(pptx_path)
        
        for i, slide in enumerate(prs.slides):
            slide_title = "Untitled Slide"
            slide_content = []
            has_charts = False
            has_images = False
            
            # Enhanced text content extraction from actual slide
            for shape in slide.shapes:
                # Check for charts in actual slide
                if hasattr(shape, 'chart'):
                    has_charts = True
                    chart_type = type(shape.chart).__name__
                    slide_content.append(f"📊 Chart: {chart_type}")
                
                # Check for images in actual slide
                if hasattr(shape, 'image'):
                    has_images = True
                    slide_content.append(f"🖼️ Image: Visual content")
                
                # Extract actual text content
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # Better title detection from actual slide
                    if shape == slide.shapes.title or (not slide_content and len(text) < 100):
                        slide_title = text
                    else:
                        # Add meaningful text content from actual slide
                        if text and text != slide_title:
                            lines = text.split('\n')
                            for line in lines:
                                if line.strip() and len(line.strip()) > 2:
                                    slide_content.append(line.strip())
            
            # Limit content for preview but show more actual content
            display_content = slide_content[:12]  # Show first 12 actual content lines
            
            # Create enhanced slide preview data with ACTUAL content
            slide_preview_data.append({
                "slide_number": i + 1,
                "title": slide_title,
                "content": display_content,
                "type": "actual_slide",
                "has_charts": has_charts,
                "has_images": has_images,
                "total_shapes": len([shape for shape in slide.shapes]),
                "is_generated": True
            })
        
        logger.info(f"Generated enhanced ACTUAL slide preview with {len(slide_preview_data)} slides")
        return slide_preview_data
        
    except Exception as e:
        logger.error(f"Error generating actual slide preview: {e}")
        return []

# API endpoints
@app.post("/upload/", response_model=FileUploadResponse)
async def upload_file(
    file: UploadFile = File(...),
    file_type: str = Form("csv")
):
    """Upload any type of file (CSV, template, logo) - NO SIZE LIMIT"""
    try:
        logger.info(f"Upload request received - Type: {file_type}, Filename: {file.filename}")
        
        # Validate file type
        if file_type not in ["csv", "template"]:
            raise HTTPException(status_code=400, detail="Invalid file type. Use 'csv' or 'template'")
        
        # Validate file extensions
        if file_type == "template" and not (file.filename or "").endswith('.pptx'):
            raise HTTPException(status_code=400, detail="Template must be PowerPoint format (.pptx)")
        
        file_id, filename = save_uploaded_file(file, file_type)
        file_path = get_file_path(file_id, file_type)
        
        return FileUploadResponse(
            file_id=file_id,
            filename=filename,
            file_type=file_type,
            file_path=str(file_path),
            message=f"{file_type.upper()} file uploaded successfully - NO SIZE LIMIT"
        )
        
    except Exception as e:
        logger.error(f"Upload error: {e}")
        raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")


@app.get("/analyze/{file_id}")
async def analyze_data(file_id: str):
    """Analyze uploaded CSV data - SUPPORTS LARGE FILES"""
    try:
        logger.info(f"Analysis request for file_id: {file_id}")
        
        # Get the actual file path
        file_path = get_file_path(file_id, "csv")
        if not file_path.exists():
            logger.error(f"File not found: {file_path}")
            raise HTTPException(status_code=404, detail="File not found")
        
        logger.info(f"Found file: {file_path}")
        
        # Use the existing analyze_csv_data function (supports large files)
        analysis = analyze_csv_data(file_path)
        
        # Convert to AnalysisResponse format
        analysis_response = AnalysisResponse(
            file_id=file_id,
            columns=analysis["columns"],
            summary=analysis["summary"],
            sample_data=analysis["sample_data"]
        )
        
        return analysis_response
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Analysis error for file {file_id}: {str(e)}")
        import traceback
        logger.error(f"Traceback: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Analysis failed: {str(e)}")

@app.post("/generate-ppt/")
async def generate_presentation(request: PPTGenerationRequest):
    """Generate PowerPoint presentation using generation pipeline - SUPPORTS LARGE FILES"""
    try:
        logger.info("PPT generation request received")
        logger.info(f"Config: {request.presentation_config}")
        
        # Validate template selection
        template_name = request.presentation_config.get("template_name", "none")
        if template_name == "none" and not request.template_file_id:
            logger.info("No template selected, using generation pipeline default")
        
        result = generate_powerpoint(request)
        
        return {
            "success": True,
            "ppt_id": result["ppt_id"],
            "total_slides": result["total_slides"],
            "slide_preview_data": result["slide_preview_data"],
            "slide_images": result.get("slide_images", []),  # ENHANCED: Include high-quality ACTUAL slide images
            "message": "Presentation generated successfully using generation pipeline",
            "download_url": f"http://localhost:8000/download/{result['ppt_id']}",
            "file_path": result["file_path"],
            "timestamp": result["generated_at"],
            "template_used": result["template_used"],
            "text_color_applied": result["text_color_applied"],
            "generation_pipeline": True
        }
        
    except Exception as e:
        logger.error(f"PPT generation error: {e}")
        raise HTTPException(status_code=500, detail=f"PPT generation failed: {str(e)}")

@app.get("/download/{ppt_id}")
async def download_presentation(ppt_id: str):
    """Download generated PowerPoint presentation"""
    try:
        file_path = OUTPUT_DIR / f"{ppt_id}.pptx"
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="Presentation not found")
        
        filename = f"data_presentation_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
        
        return FileResponse(
            path=file_path,
            filename=filename,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
    except Exception as e:
        logger.error(f"Download error: {e}")
        raise HTTPException(status_code=500, detail=f"Download failed: {str(e)}")

# ENHANCED ENDPOINT: Get high-quality ACTUAL slide images for preview
@app.get("/slide-images/{ppt_id}")
async def get_slide_images(ppt_id: str):
    """Get ACTUAL slide images from generated PowerPoint with enhanced quality and REAL CONTENT"""
    try:
        pptx_path = OUTPUT_DIR / f"{ppt_id}.pptx"
        if not pptx_path.exists():
            raise HTTPException(status_code=404, detail="Presentation not found")
        
        logger.info(f"Converting ACTUAL slides to high-quality images for: {pptx_path}")
        slide_images = convert_pptx_to_slide_images(pptx_path)
        
        if not slide_images:
            logger.warning("No slide images generated, creating fallback")
            # Create a simple fallback
            slide_images = [{
                "slide_number": 1,
                "image_data": "",
                "title": "Presentation Preview",
                "content": ["Slide images are being processed..."],
                "has_charts": False,
                "has_images": False,
                "total_shapes": 0,
                "type": "fallback"
            }]
        
        return SlideImagesResponse(
            success=True,
            ppt_id=ppt_id,
            slide_images=slide_images,
            total_slides=len(slide_images),
            timestamp=datetime.now().isoformat()
        )
        
    except Exception as e:
        logger.error(f"Error getting slide images: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting slide images: {str(e)}")

# UPDATED ENDPOINT: Get enhanced ACTUAL slide preview from generated PPTX
@app.get("/slide-preview/{csv_file_id}")
async def get_slide_preview(csv_file_id: str):
    """Get slide preview data for frontend display - UPDATED to use ACTUAL PPTX content"""
    try:
        # Get CSV file path
        csv_path = get_file_path(csv_file_id, "csv")
        if not csv_path.exists():
            raise HTTPException(status_code=404, detail="CSV file not found")
        
        # Read data for analysis
        if csv_path.suffix.lower() in ['.xlsx', '.xls', '.xlsv']:
            df = pd.read_excel(csv_path)
        else:
            df = pd.read_csv(csv_path)
        
        # Get configuration
        config_path = csv_path.parent / "input.json"
        if config_path.exists():
            with open(config_path, 'r') as f:
                config = json.load(f)
        else:
            config = {}
        
        # Check if there's a generated PPTX file we can use for ACTUAL preview
        pptx_files = list(OUTPUT_DIR.glob("*.pptx"))
        if pptx_files:
            # Use the most recent PPTX file for ACTUAL preview
            latest_pptx = max(pptx_files, key=lambda x: x.stat().st_mtime)
            ppt_id = latest_pptx.stem
            
            # Get enhanced ACTUAL slide content from PPTX
            actual_slide_preview = get_actual_slide_preview(ppt_id)
            
            if actual_slide_preview:
                logger.info(f"Using enhanced ACTUAL slide preview from {latest_pptx.name}")
                return SlidePreviewResponse(
                    success=True,
                    slide_preview_data=actual_slide_preview,
                    total_slides=len(actual_slide_preview),
                    timestamp=datetime.now().isoformat()
                )
        
        # Fallback to generated preview data if no PPTX found
        logger.info("Using generated slide preview data (no PPTX found)")
        slide_preview_data = generate_slide_preview_data(config, df)
        
        return SlidePreviewResponse(
            success=True,
            slide_preview_data=slide_preview_data,
            total_slides=len(slide_preview_data),
            timestamp=datetime.now().isoformat()
        )
        
    except Exception as e:
        logger.error(f"Error generating slide preview: {e}")
        raise HTTPException(status_code=500, detail=f"Error generating slide preview: {str(e)}")

@app.get("/health")
async def health_check():
    """Health check endpoint"""
    try:
        # Check if directories are accessible
        dir_status = {}
        for name, directory in [("input", INPUT_DIR), ("output", OUTPUT_DIR), 
                              ("csv", CSV_DIR), ("templates", TEMPLATE_DIR), ("logos", LOGO_DIR)]:
            dir_status[name] = {
                "exists": directory.exists(),
                "writable": os.access(directory, os.W_OK),
                "file_count": len(list(directory.glob("*"))) if directory.exists() else 0
            }
        
        return {
            "status": "healthy",
            "timestamp": datetime.now().isoformat(),
            "directories": dir_status,
            "version": "1.0.0"
        }
    except Exception as e:
        return {
            "status": "unhealthy",
            "error": str(e),
            "timestamp": datetime.now().isoformat()
        }

@app.post("/save-config/{csv_file_id}")
async def save_configuration(csv_file_id: str, config: dict):
    """Save presentation configuration JSON - always use input.json"""
    try:
        # Get CSV file path
        csv_path = get_file_path(csv_file_id, "csv")
        if not csv_path.exists():
            raise HTTPException(status_code=404, detail="CSV file not found")
        
        parent_dir = csv_path.parent

        # Always use input.json as filename
        config_filename = "input.json"
        config_path = parent_dir / config_filename

        # Load existing config if it exists
        if config_path.exists():
            with open(config_path, 'r') as f:
                existing_config = json.load(f)
        else:
            existing_config = {}

        # Merge the new config into the existing config
        def deep_update(orig, new):
            for key, value in new.items():
                if isinstance(value, dict) and key in orig:
                    deep_update(orig[key], value)
                else:
                    orig[key] = value

        deep_update(existing_config, config)

        # Save the updated config back (replaces existing file)
        with open(config_path, 'w') as f:
            json.dump(existing_config, f, indent=2)

        return {
            "success": True,
            "message": "Configuration updated successfully",
            "config_path": str(config_path),
            "config_filename": config_filename,
            "timestamp": datetime.now().isoformat()
        }
        
    except Exception as e:
        logger.error(f"Error saving configuration: {e}")
        raise HTTPException(status_code=500, detail=f"Error saving configuration: {str(e)}")

@app.post("/cleanup")
async def cleanup_files():
    """Manual cleanup endpoint - completely clean input and output folders"""
    try:
        success = cleanup_input_output_folders()
        
        if success:
            return {
                "success": True,
                "message": "Input and output folders cleaned successfully",
                "timestamp": datetime.now().isoformat()
            }
        else:
            raise HTTPException(status_code=500, detail="Cleanup failed")
            
    except Exception as e:
        logger.error(f"Cleanup error: {e}")
        raise HTTPException(status_code=500, detail=f"Cleanup failed: {str(e)}")

@app.get("/config/{csv_file_id}")
async def get_configuration(csv_file_id: str):
    """Get current configuration for a CSV file"""
    try:
        # Get CSV file path
        csv_path = get_file_path(csv_file_id, "csv")
        if not csv_path.exists():
            raise HTTPException(status_code=404, detail="CSV file not found")
        
        config_path = csv_path.parent / "input.json"
        
        if config_path.exists():
            with open(config_path, 'r') as f:
                config = json.load(f)
            
            return {
                "success": True,
                "config": config,
                "config_path": str(config_path),
                "timestamp": datetime.now().isoformat()
            }
        else:
            return {
                "success": True,
                "config": {},
                "message": "No configuration found",
                "timestamp": datetime.now().isoformat()
            }
            
    except Exception as e:
        logger.error(f"Error getting configuration: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting configuration: {str(e)}")

@app.get("/templates")
async def get_available_templates():
    """Get available templates"""
    try:
        templates = [
            {"name": "Professional", "value": "professional", "description": "Clean and corporate style", "icon": "💼"},
            {"name": "Creative", "value": "creative", "description": "Modern and colorful design", "icon": "🎨"},
            {"name": "Minimal", "value": "minimal", "description": "Simple and elegant layout", "icon": "⚪"},
            {"name": "Technical", "value": "technical", "description": "Data-focused with charts emphasis", "icon": "📈"}
        ]
        
        return TemplateResponse(
            success=True,
            templates=templates,
            current_template="professional"
        )
        
    except Exception as e:
        logger.error(f"Error getting templates: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting templates: {str(e)}")

@app.get("/files/list")
async def list_uploaded_files():
    """List all uploaded files for debugging"""
    try:
        files_info = {}
        
        for dir_name, directory in [("csv", CSV_DIR), ("templates", TEMPLATE_DIR), ("logos", LOGO_DIR), ("output", OUTPUT_DIR)]:
            if directory.exists():
                files = []
                for file_path in directory.glob("*"):
                    if file_path.is_file():
                        file_size = file_path.stat().st_size / (1024 * 1024)  # Size in MB
                        files.append({
                            "name": file_path.name,
                            "size_mb": round(file_size, 2),
                            "modified": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat()
                        })
                files_info[dir_name] = files
            else:
                files_info[dir_name] = "Directory not found"
        
        return {
            "success": True,
            "files": files_info,
            "timestamp": datetime.now().isoformat()
        }
        
    except Exception as e:
        logger.error(f"Error listing files: {e}")
        raise HTTPException(status_code=500, detail=f"Error listing files: {str(e)}")

# Background task for cleanup
async def cleanup_old_files():
    """Clean up files older than 1 hour (more aggressive cleanup)"""
    while True:
        try:
            current_time = datetime.now().timestamp()
            max_age = 1 * 60 * 60  # 1 hour in seconds (more aggressive)
            
            # Clean up output directory files older than 1 hour
            if OUTPUT_DIR.exists():
                for file_path in OUTPUT_DIR.glob("*.pptx"):
                    if file_path.is_file() and (current_time - file_path.stat().st_mtime > max_age):
                        file_path.unlink()
                        logger.info(f"Cleaned up old presentation: {file_path.name}")
            
            # Clean up old template files (24 hours)
            max_age_24h = 24 * 60 * 60
            for directory in [TEMPLATE_DIR]:
                if directory.exists():
                    for file_path in directory.glob("*"):
                        if file_path.is_file() and (current_time - file_path.stat().st_mtime > max_age_24h):
                            file_path.unlink()
                            logger.info(f"Cleaned up old file: {file_path.name}")
                        
            await asyncio.sleep(1800)  # Run every 30 minutes
            
        except Exception as e:
            logger.error(f"Cleanup error: {e}")
            await asyncio.sleep(300)  # Wait 5 minutes on error

@app.on_event("startup")
async def startup_event():
    """Start background tasks on startup"""
    logger.info("Starting Excel to PPT Generator API")
    logger.info(f"Input directory: {INPUT_DIR}")
    logger.info(f"Output directory: {OUTPUT_DIR}")
    
    # Ensure directories exist
    for directory in [INPUT_DIR, OUTPUT_DIR, CSV_DIR, TEMPLATE_DIR, LOGO_DIR]:
        directory.mkdir(parents=True, exist_ok=True)
    
    # Start cleanup task
    asyncio.create_task(cleanup_old_files())

@app.on_event("shutdown")
async def shutdown_event():
    """Cleanup on shutdown"""
    logger.info("Shutting down Excel to PPT Generator API")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000, log_level="info")
    
